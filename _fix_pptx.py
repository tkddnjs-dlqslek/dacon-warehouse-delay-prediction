import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from lxml import etree

pptx_path = r'C:\Users\user\Desktop\데이콘 4월\dacon_smart_warehouse.pptx'
prs = Presentation(pptx_path)

slide = prs.slides[5]  # p.6 모델 파이프라인

# shape[24] = scenario_id 잘못된 줄 확인
target_shape = slide.shapes[24]
print("삭제 대상:", repr(target_shape.text_frame.paragraphs[0].text))

# XML에서 shape 제거
sp = target_shape._element
sp.getparent().remove(sp)

# shape[22] (layout_id 올바른 줄) 확인
correct_shape = slide.shapes[22]
print("유지 대상:", repr(correct_shape.text_frame.paragraphs[0].text))

# 저장
prs.save(pptx_path)
print("PPTX 저장 완료")
